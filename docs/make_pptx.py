"""Generate esitlus.pptx for Hoonete soojakadu project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1B, 0x37, 0x6B)   # headings / title bg
MID_BLUE    = RGBColor(0x2E, 0x75, 0xB6)   # accent
LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)   # table header bg
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x40, 0x40, 0x40)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE      = RGBColor(0xED, 0x7D, 0x31)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=DARK_GREY,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_filled_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def bullet_frame(slide, items, left, top, width, height,
                 font_size=16, color=DARK_GREY, bullet="•"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    import copy
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, DARK_BLUE)

# Decorative orange bar
add_filled_rect(s1, 0, 3.0, 13.33, 0.07, ORANGE)

add_textbox(s1, "Hoonete soojakadu ja renoveerimisprioriteedid",
            0.6, 1.0, 12.0, 1.8,
            font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s1, "Andmeinseneeria kursuse grupitöö  ·  2025",
            0.6, 2.9, 12.0, 0.6,
            font_size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_textbox(s1,
            "Anna-Liisa Altmets  ·  Erkki Laaneoks  ·  Mihkel Nugis  ·  Toel Teemaa",
            0.6, 3.4, 12.0, 0.6,
            font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
add_textbox(s1,
            "Millistesse omavalitsustesse läheb kõige rohkem soojusenergiat raisku\n"
            "ning kus annab renoveerimine kõige suurema efekti?",
            1.2, 4.5, 10.8, 1.4,
            font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Probleem ja andmeallikad
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, WHITE)

add_filled_rect(s2, 0, 0, 13.33, 1.1, DARK_BLUE)
add_textbox(s2, "Probleem ja andmeallikad",
            0.4, 0.1, 12.5, 0.85,
            font_size=28, bold=True, color=WHITE)

# Left column — problem
add_filled_rect(s2, 0.3, 1.25, 5.9, 5.6, LIGHT_GREY)
add_textbox(s2, "Miks see oluline on?",
            0.5, 1.35, 5.5, 0.5,
            font_size=18, bold=True, color=DARK_BLUE)
bullet_frame(s2, [
    "Eesti hoonefondist suur osa ehitatud nõukogude ajal — halb soojustus",
    "Soojus läheb läbi seinte ja katuste kaduma → kõrged energiarved + CO₂",
    "Riiklikud renoveerimiseesmärgid nõuavad andmepõhist prioritiseerimist",
    "Küsimus: kus oleks renoveerimine kõige suurema mõjuga?",
], 0.5, 1.95, 5.6, 4.5, font_size=15)

# Right column — data sources
add_filled_rect(s2, 7.0, 1.25, 5.9, 5.6, LIGHT_GREY)
add_textbox(s2, "Kolm andmeallikat",
            7.2, 1.35, 5.5, 0.5,
            font_size=18, bold=True, color=DARK_BLUE)

sources = [
    ("Ehitisregister", "~300 000 hoonet\nehitusaasta · pindala · asukoht"),
    ("Open-Meteo API", "Päevased keskmised temperatuurid\n6 ilmajaamast üle Eesti"),
    ("Elering / Nord Pool", "Tunnihinnad börsil (€/MWh)\nTegelik tarbimine (MWh)"),
]
y = 2.0
for name, desc in sources:
    add_filled_rect(s2, 7.15, y, 5.55, 1.35, MID_BLUE)
    add_textbox(s2, name, 7.3, y + 0.05, 5.2, 0.45,
                font_size=15, bold=True, color=WHITE)
    add_textbox(s2, desc, 7.3, y + 0.5, 5.2, 0.75,
                font_size=13, color=WHITE)
    y += 1.55

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Metoodika: soojakao valem
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, WHITE)

add_filled_rect(s3, 0, 0, 13.33, 1.1, DARK_BLUE)
add_textbox(s3, "Metoodika: soojakao arvutus",
            0.4, 0.1, 12.5, 0.85,
            font_size=28, bold=True, color=WHITE)

# Formula box
add_filled_rect(s3, 1.5, 1.3, 10.3, 1.3, LIGHT_BLUE)
add_textbox(s3,
            "Soojakadu (kWh)  =  U-väärtus  ×  Pindala  ×  (21°C − Õhutemperatuur)  ×  24h",
            1.7, 1.45, 9.9, 1.0,
            font_size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Explanations in 3 boxes
boxes = [
    ("U-väärtus (W/m²K)",
     "Näitab soojapidavust.\nVanem hoone = suurem U = rohkem kadu.\nLookup-tabel ehitusaasta järgi."),
    ("Temperatuurivahe",
     "Sisetemperatuur 21 °C miinus\ntegelik välisõhu temp.\nKülmem ilm = suurem kadu."),
    ("Agregeerib iga päev",
     "Arvutus iga hoone × iga päev.\nSumma omavalitsuse tasemele.\n~110 mln rida aastas."),
]
x = 0.5
for title, body in boxes:
    add_filled_rect(s3, x, 2.9, 3.9, 2.8, LIGHT_GREY)
    add_textbox(s3, title, x + 0.15, 2.97, 3.6, 0.55,
                font_size=15, bold=True, color=MID_BLUE)
    add_textbox(s3, body, x + 0.15, 3.55, 3.6, 2.0,
                font_size=13, color=DARK_GREY)
    x += 4.15

add_textbox(s3,
            "Tulemus: neli mõõdikut iga omavalitsuse kohta — "
            "intensiivsus · ilmastikutundlikkus · renoveerimispotentsiaal · mudeli valideerimine",
            0.5, 6.15, 12.3, 0.9,
            font_size=14, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tehniline arhitektuur
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, WHITE)

add_filled_rect(s4, 0, 0, 13.33, 1.1, DARK_BLUE)
add_textbox(s4, "Tehniline arhitektuur",
            0.4, 0.1, 12.5, 0.85,
            font_size=28, bold=True, color=WHITE)

# Pipeline arrow row
pipeline = [
    ("Andmeallikad", "Ehitisregister\nOpen-Meteo\nElering"),
    ("Airflow", "Iga-päevane\nDAG"),
    ("PostgreSQL", "Toortabelid\n(UPSERT)"),
    ("dbt", "Staging →\nIntermediate\n→ Marts"),
    ("Superset", "Interaktiivne\nkaart ja\ngraafikud"),
]
colors = [DARK_BLUE, MID_BLUE, MID_BLUE, MID_BLUE, DARK_BLUE]
x = 0.25
for i, (title, sub) in enumerate(pipeline):
    add_filled_rect(s4, x, 1.4, 2.3, 1.8, colors[i])
    add_textbox(s4, title, x + 0.1, 1.45, 2.1, 0.55,
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s4, sub, x + 0.05, 2.0, 2.2, 1.1,
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        add_textbox(s4, "→", x + 2.3, 1.9, 0.45, 0.6,
                    font_size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    x += 2.58

# dbt layers detail
add_textbox(s4, "dbt kihid",
            0.4, 3.55, 5.0, 0.5,
            font_size=16, bold=True, color=DARK_BLUE)
layers = [
    ("Staging", "Andmete puhastamine ja normaliseerimine"),
    ("Intermediate", "Soojakao arvutus (~110 mln rida/aastas)"),
    ("Marts", "Lõplikud analüütilised tabelid Superseti jaoks"),
]
y = 4.1
for layer, desc in layers:
    add_filled_rect(s4, 0.4, y, 1.6, 0.5, MID_BLUE)
    add_textbox(s4, layer, 0.45, y + 0.03, 1.5, 0.44,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s4, desc, 2.2, y + 0.05, 4.5, 0.44,
                font_size=13, color=DARK_GREY)
    y += 0.65

# Right side — quality & infra
add_textbox(s4, "Kvaliteet ja infrastruktuur",
            7.2, 3.55, 5.8, 0.5,
            font_size=16, bold=True, color=DARK_BLUE)
bullet_frame(s4, [
    "dbt test — automaatsed andmekvaliteedi testid peale iga laadimist",
    "dbt snapshot — omavalitsuste ühinemiste jälgimine (SCD Type 2)",
    "Docker Compose — kogu keskkond käivitub ühe käsuga",
    "GitHub — versioonikontroll, meeskonna koostöö",
], 7.2, 4.1, 5.8, 3.0, font_size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Tulemused ja mõõdikud
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, WHITE)

add_filled_rect(s5, 0, 0, 13.33, 1.1, DARK_BLUE)
add_textbox(s5, "Tulemused: neli mõõdikut omavalitsuste kaupa",
            0.4, 0.1, 12.5, 0.85,
            font_size=28, bold=True, color=WHITE)

metrics = [
    ("Soojakao intensiivsus",
     "kWh/m²/aastas",
     "Mida suurem, seda halvem soojapidavus.\nPiirkondlik pingerida — kus on kõige külmemad hooned?",
     DARK_BLUE),
    ("Ilmastikutundlikkus",
     "Pearsoni r",
     "Kui tugevalt korreleerub soojakadu\nvälistemperatuuriga? (–1 kuni +1)",
     MID_BLUE),
    ("Renoveerimispotentsiaal",
     "€ / aastas",
     "Potentsiaalne aastane kokkuhoid eurodes,\nkui viia intensiivsus 50 kWh/m²-ni.",
     ORANGE),
    ("Mudeli valideerimine",
     "Pearsoni r + KAV",
     "Mudeli soojakadu vs. tegelik Eleringi\ntarbimine — kui täpne mudel on?",
     rgb(0x5A, 0x9E, 0x6F)),
]

x = 0.3
for title, unit, desc, col in metrics:
    add_filled_rect(s5, x, 1.3, 3.0, 0.7, col)
    add_textbox(s5, title, x + 0.1, 1.33, 2.8, 0.65,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(s5, x, 2.0, 3.0, 0.5, LIGHT_BLUE)
    add_textbox(s5, unit, x + 0.05, 2.03, 2.9, 0.44,
                font_size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s5, desc, x + 0.1, 2.6, 2.8, 1.5,
                font_size=12, color=DARK_GREY)
    x += 3.2

# Superset note
add_filled_rect(s5, 0.5, 4.4, 12.3, 2.5, LIGHT_GREY)
add_textbox(s5, "Visualiseerimine — Apache Superset",
            0.7, 4.5, 8.0, 0.5,
            font_size=16, bold=True, color=DARK_BLUE)
bullet_frame(s5, [
    "Interaktiivne kaart: soojakao intensiivsus Eesti omavalitsuste kaupa (choropleth)",
    "Tulpdiagramm: Top-20 omavalitsust renoveerimispotentsiaali järgi (€/aastas)",
    "Hajuvusdiagramm: ilmastikutundlikkus vs soojakao intensiivsus",
    "Ajagraafik: mudeli vs tegeliku tarbimise võrdlus",
], 0.7, 5.05, 12.0, 1.75, font_size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Tööjaotus ja kokkuvõte
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, WHITE)

add_filled_rect(s6, 0, 0, 13.33, 1.1, DARK_BLUE)
add_textbox(s6, "Tööjaotus ja kokkuvõte",
            0.4, 0.1, 12.5, 0.85,
            font_size=28, bold=True, color=WHITE)

# Team roles
team = [
    ("Toel Teemaa",    "Andmete laadimine",  "Airflow DAG-id, API integratsioonid (Open-Meteo, Elering, Nord Pool)"),
    ("Erkki Laaneoks", "Teisendused",         "dbt mudelid, soojakao arvutus, mart-kiht, pearson_r makro"),
    ("Mihkel Nugis",   "Andmekvaliteet",      "dbt testid, vigade diagnoos, andmete valideerimise automatiseerimine"),
    ("Anna-Liisa Altmets", "Visualiseerimine", "Superset näidikulaud, äriküsimuse sidumine tulemustega"),
]
cols = [MID_BLUE, DARK_BLUE, rgb(0x5A, 0x9E, 0x6F), ORANGE]
x = 0.25
for i, (name, role, desc) in enumerate(team):
    add_filled_rect(s6, x, 1.3, 2.95, 0.65, cols[i])
    add_textbox(s6, name, x + 0.1, 1.32, 2.75, 0.6,
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_filled_rect(s6, x, 1.95, 2.95, 0.48, LIGHT_BLUE)
    add_textbox(s6, role, x + 0.05, 1.97, 2.85, 0.44,
                font_size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s6, desc, x + 0.1, 2.5, 2.78, 1.4,
                font_size=12, color=DARK_GREY)
    x += 3.18

# Summary
add_filled_rect(s6, 0.4, 4.2, 12.5, 2.95, LIGHT_GREY)
add_textbox(s6, "Kokkuvõte",
            0.6, 4.3, 5.0, 0.5,
            font_size=16, bold=True, color=DARK_BLUE)
bullet_frame(s6, [
    "Ühendame 300 000 hoone andmed ilmaandmete ja energiahindadega → omavalitsuste pingerida",
    "Automatiseeritud torujuhe: Airflow + dbt + PostgreSQL + Superset, kõik Dockeris",
    "Andmekvaliteedi testid tagavad tulemuste usaldusväärsuse (automaatne kontroll iga laadimise järel)",
    "Tulemus on praktiline tööriist renoveerimiseelarve suunamiseks seal, kus mõju on suurim",
], 0.6, 4.85, 8.5, 2.2, font_size=14)

# Stack logos text
add_textbox(s6, "Stack:",
            9.3, 4.3, 1.0, 0.4,
            font_size=13, bold=True, color=DARK_BLUE)
stack_items = ["Apache Airflow", "dbt Core", "PostgreSQL", "Apache Superset", "Docker"]
y = 4.75
for item in stack_items:
    add_filled_rect(s6, 9.3, y, 3.2, 0.38, MID_BLUE)
    add_textbox(s6, item, 9.35, y + 0.02, 3.1, 0.34,
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    y += 0.46

# Save
out = r"c:\Users\mihkel\OneDrive - Digmatix\Documents\Repos\Andmeinseneeria\Grupitöö\docs\esitlus.pptx"
prs.save(out)
print(f"Saved: {out}")
