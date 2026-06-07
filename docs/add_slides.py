"""Append dashboard result slides to existing esitlus.pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BLUE  = RGBColor(0x1B, 0x37, 0x6B)
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x40, 0x40, 0x40)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)
GREEN      = RGBColor(0x5A, 0x9E, 0x6F)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, left, top, width, height,
       size=14, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def rect(slide, left, top, width, height, fill, line=None):
    sh = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh


def bullets(slide, items, left, top, width, height, size=13, color=DARK_GREY):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(3)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color


PPTX = r"c:\Users\mihkel\OneDrive - Digmatix\Documents\Repos\Andmeinseneeria\Grupitöö\docs\esitlus.pptx"
prs = Presentation(PPTX)
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# SLAID 7 — Dashboard: põhitulemused
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
set_bg(s7, WHITE)

rect(s7, 0, 0, 13.33, 1.1, DARK_BLUE)
tb(s7, "Dashboard: põhitulemused", 0.4, 0.1, 12.5, 0.85,
   size=28, bold=True, color=WHITE)

# 4 KPI kaarti
kpis = [
    ("Soojakadu Eestis", "9.26B kWh",    DARK_BLUE),
    ("Kulu",             "660M €",        MID_BLUE),
    ("Säästupotentsiaal","423M €/aastas", GREEN),
    ("Keskmine intensiivsus", "149.93 kWh/m²", ORANGE),
]
x = 0.3
for label, value, col in kpis:
    rect(s7, x, 1.25, 2.95, 1.5, col)
    tb(s7, label, x+0.1, 1.3, 2.75, 0.55, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s7, value, x+0.05, 1.85, 2.85, 0.75, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += 3.18

# Graafikute kokkuvõtted
charts = [
    ("Suurim vajadus renoveerimiseks (vajaduse skoor)",
     "Lääneranna vald · Kihnu vald · Räpina vald · Toila vald · Kanepi vald\n"
     "Mõõdik: soojakadu_kwh_aastas / kogupindala_m²"),
    ("TOP 10 KOV — suurim aastane säästupotentsiaal (€)",
     "Tallinn ~100M € · Tartu linn ~20M € · Pärnu linn · Narva linn\n"
     "Kohtla-Järve · Saaremaa vald · Viljandi vald · Jõgeva vald · Elva vald"),
    ("Investeeringu tasuvus — ROI vs tasuvusaeg",
     "Kambja vald (ROI 3.4) · Saue vald (3.36) · Tori vald (3.31)\n"
     "Rapla vald (3.21) — finantsiliselt kõige tasuvam renoveerimine"),
]
x = 0.3
for title, body in charts:
    rect(s7, x, 3.0, 4.15, 3.8, LIGHT_GREY)
    tb(s7, title, x+0.12, 3.07, 3.9, 0.65, size=12, bold=True, color=MID_BLUE)
    tb(s7, body,  x+0.12, 3.75, 3.9, 2.9, size=12, color=DARK_GREY)
    x += 4.35

# ══════════════════════════════════════════════════════════════════════════════
# SLAID 8 — Järeldused ja ärimeetrikud
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
set_bg(s8, WHITE)

rect(s8, 0, 0, 13.33, 1.1, DARK_BLUE)
tb(s8, "Järeldused ja ärimeetrikud", 0.4, 0.1, 12.5, 0.85,
   size=28, bold=True, color=WHITE)

# Järeldused
rect(s8, 0.3, 1.25, 7.8, 4.75, LIGHT_GREY)
tb(s8, "Peamised järeldused", 0.5, 1.32, 7.4, 0.5,
   size=16, bold=True, color=DARK_BLUE)
bullets(s8, [
    "Finantsiliselt tasuvam: Kambja vald (ROI 3.4) ja Saue vald (ROI 3.36)",
    "Suurim absoluutne säästupotentsiaal: Tallinn (üle 100M €/aastas)",
    "Vajaduspõhiselt haavatavaimad: Lääneranna, Kihnu ja Räpina vald",
    "Tallinna ja Tartu linna ROI on madalam (2.26–2.27), aga kogu mõju suurim",
    "Kokkuvõte: kõrge ROI → väiksemad vallad; suur säästusumma → linnad",
], 0.5, 1.9, 7.5, 3.9, size=14)

# Ärimeetrikud
rect(s8, 8.4, 1.25, 4.6, 4.75, DARK_BLUE)
tb(s8, "Ärimeetrikute valemid", 8.55, 1.32, 4.3, 0.5,
   size=14, bold=True, color=WHITE)

formulas = [
    ("ROI",             "sääst_eur_aastas\n/ investeering"),
    ("VAJADUSE SKOOR",  "soojakadu_kwh\n/ kogupindala_m²"),
    ("INVESTEERING",    "kogupindala_m² × 250\n(€/m² proxy)"),
    ("PRIORITY",        "LOG(sääst+1) × LOG(skoor+1)\n× LOG(pindala+1)"),
    ("TASUVUS",         "1 / NULLIF(roi, 0)"),
]
y = 1.95
for name, formula in formulas:
    rect(s8, 8.45, y, 4.45, 0.82, MID_BLUE)
    tb(s8, name,    8.57, y+0.02, 4.2, 0.35, size=11, bold=True, color=WHITE)
    tb(s8, formula, 8.57, y+0.38, 4.2, 0.42, size=10, color=LIGHT_BLUE)
    y += 0.9

prs.save(PPTX)
print(f"Lisatud 2 slaidi: {PPTX}")
